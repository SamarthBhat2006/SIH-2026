"""
AI PowerPoint (.pptx) Generator Module
Isolated module for transforming source content into professional PowerPoint presentations
using Google Gemini API (google-genai) and python-pptx.
"""

import os
import json
import time
import re
from pathlib import Path
from typing import Dict, Any, List, Optional
from dotenv import load_dotenv

# Ensure environment variables are loaded
load_dotenv()

class PPTGenerationError(Exception):
    """Custom exception for user-friendly PPT generation error messaging."""
    pass

def _clean_json_response(raw_text: str) -> str:
    """Removes Markdown code blocks or wrapping text to extract valid JSON string."""
    cleaned = raw_text.strip()
    if cleaned.startswith("```"):
        # Strip ```json ... ``` fences
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    return cleaned.strip()

def generate_ppt_content(
    source_text: str,
    audience: str = "Executive",
    tone: str = "Professional"
) -> Dict[str, Any]:
    """
    Transforms source content into a structured 5-7 slide presentation JSON using Google Gemini API.
    Uses ONLY Gemini API with strict grounding instructions.
    
    Returns parsed presentation dictionary.
    Raises PPTGenerationError with exact user-facing error messages on failure.
    """
    # 1. Verify GEMINI_API_KEY configuration
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not gemini_key:
        raise PPTGenerationError(
            "Gemini API key is not configured. Please configure GEMINI_API_KEY in the .env file."
        )

    # 2. Build strict grounding prompt
    system_instruction = (
        "You are an expert executive presentation planner for high-consequence technical briefings. "
        "Your task is to convert the provided source content into a professional 5 to 7 slide presentation structure.\n\n"
        "STRICT GROUNDING RULES:\n"
        "1. Use ONLY information explicitly contained in the provided source text.\n"
        "2. NEVER invent facts, statistics, entities, or timeline events.\n"
        "3. Clearly indicate when requested information is unavailable in the source.\n"
        "4. Create concise slide content suitable for presentation slides.\n"
        "5. Preserve important technical and cybersecurity terminology.\n"
        "6. Include informative speaker notes for each slide.\n"
        "7. Structure the presentation logically (e.g., Title, Overview, Key Findings, Details/Evidence, Impact/Risk, Recommendations, Conclusion).\n"
        "8. Output MUST be VALID JSON ONLY. Do NOT include markdown formatting or ```json code block fences."
    )

    prompt = f"""Target Audience: {audience}
Tone: {tone}

SOURCE CONTENT:
\"\"\"
{source_text}
\"\"\"

Return a JSON object with this exact structure:
{{
  "presentation_title": "Concise Main Title",
  "subtitle": "Subtitle describing the brief for {audience}",
  "slides": [
    {{
      "slide_number": 1,
      "title": "Title",
      "layout": "title",
      "bullets": [],
      "speaker_notes": "Speaker notes for slide 1..."
    }},
    {{
      "slide_number": 2,
      "title": "Executive Overview",
      "layout": "content",
      "bullets": [
        "Concise bullet point 1 strictly from source",
        "Concise bullet point 2 strictly from source",
        "Concise bullet point 3 strictly from source"
      ],
      "speaker_notes": "Speaker notes explaining the overview..."
    }}
  ]
}}

Allowed layout values: "title", "content", "two_column", "conclusion".
Return ONLY the raw JSON string without any markdown markdown wrapper or backticks.
"""

    # 3. Invoke Google Gemini API via official google-genai SDK
    try:
        from google import genai
        from google.genai import types
        
        client = genai.Client(api_key=gemini_key)
        
        models_to_try = [
            "gemini-2.5-flash",
            "gemini-3.6-flash",
        ]
        
        config = types.GenerateContentConfig(
            system_instruction=system_instruction,
            temperature=0.2,
        )
        
        response_text = None
        last_error = None
        
        for model_name in models_to_try:
            try:
                res = client.models.generate_content(
                    model=model_name,
                    contents=prompt,
                    config=config
                )
                if res and res.text:
                    response_text = res.text
                    break
            except Exception as e:
                last_error = e
                continue
                
        if not response_text:
            raise PPTGenerationError("Unable to generate presentation content. Please try again.")

    except PPTGenerationError:
        raise
    except Exception as e:
        raise PPTGenerationError("Unable to generate presentation content. Please try again.")

    # 4. Parse and Validate JSON structure
    try:
        cleaned_text = _clean_json_response(response_text)
        data = json.loads(cleaned_text)
        
        if not isinstance(data, dict) or "slides" not in data or not isinstance(data["slides"], list):
            raise PPTGenerationError("AI returned an invalid presentation structure. Please try again.")
            
        if len(data["slides"]) == 0:
            raise PPTGenerationError("AI returned an invalid presentation structure. Please try again.")
            
        return data
        
    except json.JSONDecodeError:
        raise PPTGenerationError("AI returned an invalid presentation structure. Please try again.")
    except PPTGenerationError:
        raise
    except Exception:
        raise PPTGenerationError("AI returned an invalid presentation structure. Please try again.")

def build_pptx_file(
    presentation_data: Dict[str, Any],
    output_dir: Optional[Path] = None
) -> Path:
    """
    Converts structured presentation JSON data into a styled .pptx PowerPoint file.
    
    Returns absolute Path to the generated .pptx file.
    Raises PPTGenerationError if generation fails.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        raise PPTGenerationError("Unable to create the PowerPoint file.")

    try:
        if output_dir is None:
            output_dir = Path(__file__).resolve().parent.parent / "outputs"
        output_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        # Set Widescreen 16:9 Aspect Ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Color Palette Definitions
        COLOR_PRIMARY = RGBColor(10, 25, 47)      # Deep Navy Header
        COLOR_ACCENT = RGBColor(0, 102, 204)      # Accent Blue
        COLOR_BG_CARD = RGBColor(248, 250, 252)   # Soft Light Slate
        COLOR_TEXT_MAIN = RGBColor(30, 41, 59)    # Dark Charcoal
        COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Muted Gray
        COLOR_WHITE = RGBColor(255, 255, 255)

        title = presentation_data.get("presentation_title", "Operational Presentation")
        subtitle = presentation_data.get("subtitle", "Grounded Technical Briefing")
        slides_list = presentation_data.get("slides", [])

        total_slides = len(slides_list)

        for idx, s_data in enumerate(slides_list):
            s_num = s_data.get("slide_number", idx + 1)
            s_title = s_data.get("title", f"Slide {s_num}")
            s_layout = s_data.get("layout", "content").lower()
            bullets = s_data.get("bullets", [])
            notes = s_data.get("speaker_notes", "")

            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # A. Title Slide Layout
            if idx == 0 or s_layout == "title":
                # Dark Header Accent Card
                top_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
                )
                top_bg.fill.solid()
                top_bg.fill.fore_color.rgb = COLOR_PRIMARY
                top_bg.line.fill.background()

                # Inner Accent Banner Line
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.3), Inches(10.933), Inches(0.06)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = COLOR_ACCENT
                line.line.fill.background()

                # Title Text Box
                tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(2.2))
                tf = tb.text_frame
                tf.word_wrap = True
                p_title = tf.paragraphs[0]
                p_title.text = s_title if s_title and s_title != f"Slide {s_num}" else title
                p_title.font.size = Pt(36)
                p_title.font.bold = True
                p_title.font.name = "Calibri"
                p_title.font.color.rgb = COLOR_WHITE

                # Subtitle Text Box
                tb_sub = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.933), Inches(1.5))
                tf_sub = tb_sub.text_frame
                tf_sub.word_wrap = True
                p_sub = tf_sub.paragraphs[0]
                p_sub.text = subtitle
                p_sub.font.size = Pt(20)
                p_sub.font.name = "Calibri"
                p_sub.font.color.rgb = RGBColor(203, 213, 225)

            # B. Content / Two-Column / Conclusion Layouts
            else:
                # Top Header Banner
                header = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
                )
                header.fill.solid()
                header.fill.fore_color.rgb = COLOR_PRIMARY
                header.line.fill.background()

                # Slide Title
                tb_h = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.8))
                tf_h = tb_h.text_frame
                tf_h.word_wrap = True
                p_h = tf_h.paragraphs[0]
                p_h.text = s_title
                p_h.font.size = Pt(24)
                p_h.font.bold = True
                p_h.font.name = "Calibri"
                p_h.font.color.rgb = COLOR_WHITE

                # Two Column Layout
                if s_layout == "two_column" and len(bullets) >= 2:
                    half = (len(bullets) + 1) // 2
                    left_bullets = bullets[:half]
                    right_bullets = bullets[half:]

                    # Left Box
                    tb_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0))
                    tf_l = tb_l.text_frame
                    tf_l.word_wrap = True
                    for b_idx, b_text in enumerate(left_bullets):
                        p = tf_l.paragraphs[0] if b_idx == 0 else tf_l.add_paragraph()
                        p.text = f"•  {b_text}"
                        p.font.size = Pt(16)
                        p.font.name = "Calibri"
                        p.font.color.rgb = COLOR_TEXT_MAIN
                        p.space_after = Pt(12)

                    # Right Box
                    tb_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.0))
                    tf_r = tb_r.text_frame
                    tf_r.word_wrap = True
                    for b_idx, b_text in enumerate(right_bullets):
                        p = tf_r.paragraphs[0] if b_idx == 0 else tf_r.add_paragraph()
                        p.text = f"•  {b_text}"
                        p.font.size = Pt(16)
                        p.font.name = "Calibri"
                        p.font.color.rgb = COLOR_TEXT_MAIN
                        p.space_after = Pt(12)

                # Standard Content / Conclusion Layout
                else:
                    tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(4.9))
                    tf_b = tb_b.text_frame
                    tf_b.word_wrap = True
                    
                    if not bullets:
                        bullets = ["(No additional bullet detail specified in source content)"]

                    for b_idx, b_text in enumerate(bullets):
                        p = tf_b.paragraphs[0] if b_idx == 0 else tf_b.add_paragraph()
                        p.text = f"•  {b_text}"
                        p.font.size = Pt(17) if s_layout == "conclusion" else Pt(16)
                        p.font.name = "Calibri"
                        p.font.bold = (s_layout == "conclusion")
                        p.font.color.rgb = COLOR_TEXT_MAIN
                        p.space_after = Pt(14)

                # Bottom Footer Bar
                tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.0), Inches(0.4))
                tf_f = tb_f.text_frame
                p_f = tf_f.paragraphs[0]
                p_f.text = f"NTRO AI Platform  |  {title}"
                p_f.font.size = Pt(10)
                p_f.font.color.rgb = COLOR_TEXT_MUTED

                # Slide Number Box
                tb_num = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.0), Inches(0.4))
                tf_num = tb_num.text_frame
                p_num = tf_num.paragraphs[0]
                p_num.alignment = PP_ALIGN.RIGHT
                p_num.text = f"Slide {idx + 1} of {total_slides}"
                p_num.font.size = Pt(10)
                p_num.font.color.rgb = COLOR_TEXT_MUTED

            # Add Speaker Notes safely if supported
            if notes:
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_tf = slide.notes_slide.notes_text_frame
                        notes_tf.text = notes
                except Exception:
                    pass

        # Save PPTX with unique timestamp filename
        timestamp = int(time.time())
        filename = f"AI_Generated_Presentation_{timestamp}.pptx"
        output_path = output_dir / filename
        prs.save(str(output_path))
        return output_path

    except PPTGenerationError:
        raise
    except Exception as e:
        raise PPTGenerationError("Unable to create the PowerPoint file.")
